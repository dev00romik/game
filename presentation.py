from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ===== Налаштування презентації =====
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== Палітра кольорів =====
BG_COLOR = RGBColor(0xFD, 0xF6, 0xF0)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR = RGBColor(0x3D, 0x2C, 0x2A)
ACCENT = RGBColor(0xC7, 0x7D, 0x7D)
PHENOM_BG = RGBColor(0xFD, 0xF2, 0xF2)
PHENOM_BORDER = RGBColor(0xF5, 0xE0, 0xE0)
PHENOM_TEXT = RGBColor(0xB8, 0x65, 0x65)
FORMULA_BG = RGBColor(0xFD, 0xF6, 0xF0)
FORMULA_BORDER = RGBColor(0xD4, 0xA5, 0xA5)
FORMULA_TEXT = RGBColor(0x7A, 0x4A, 0x4A)
LABEL_COLOR = RGBColor(0x5C, 0x3D, 0x3D)
BODY_COLOR = RGBColor(0x4A, 0x35, 0x35)
SUBTITLE_COLOR = RGBColor(0x8C, 0x6E, 0x6E)
MOTTO_COLOR = RGBColor(0xB0, 0x8A, 0x8A)
EXTRA_COLOR = RGBColor(0x7A, 0x5E, 0x5E)
MORNING_BG = RGBColor(0xFE, 0xF3, 0xE4)
MORNING_TEXT = RGBColor(0xC0, 0x78, 0x40)
DAY_BG = RGBColor(0xFE, 0xF9, 0xF3)
DAY_TEXT = RGBColor(0xB0, 0x8A, 0x50)
EVENING_BG = RGBColor(0xF5, 0xF0, 0xFA)
EVENING_TEXT = RGBColor(0x7B, 0x5E, 0x9C)
CARD_BORDER = RGBColor(0xE8, 0xE0, 0xE0)

# ===== Розміри картки =====
CARD_LEFT = Inches(0.6)
CARD_TOP = Inches(0.4)
CARD_WIDTH = Inches(12.133)
CARD_HEIGHT = Inches(6.7)

# ===== Відступи всередині картки =====
LEFT = Inches(1.0)
CONTENT_WIDTH = Inches(11.133)
BADGE_TOP = Inches(0.9)
BADGE_HEIGHT = Inches(0.38)
TITLE_TOP = Inches(1.5)
TITLE_HEIGHT = Inches(0.6)
PHENOMENON_TOP = Inches(2.2)
PHENOMENON_HEIGHT = Inches(0.38)
TEXT_START_TOP = Inches(2.9)
PARAGRAPH_HEIGHT = Inches(0.28)
FORMULA_HEIGHT = Inches(0.42)
BOLD_HEIGHT = Inches(0.34)
SPACING = Inches(0.08)

# ===== Допоміжні функції =====
def set_slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_card(slide):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, CARD_LEFT, CARD_TOP, CARD_WIDTH, CARD_HEIGHT
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1)
    return card

def add_badge(slide, text, bg_color, text_color, left=LEFT, top=BADGE_TOP):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.4), BADGE_HEIGHT
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    return badge

def add_title(slide, text, left=LEFT, top=TITLE_TOP):
    txBox = slide.shapes.add_textbox(left, top, Inches(11.133), TITLE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    return txBox

def add_phenomenon(slide, text, left=LEFT, top=PHENOMENON_TOP):
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.5), PHENOMENON_HEIGHT
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = PHENOM_BG
    tag.line.color.rgb = PHENOM_BORDER
    tag.line.width = Pt(1)
    p = tag.text_frame.paragraphs[0]
    p.text = "⚡ " + text
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PHENOM_TEXT
    return tag

def add_textbox(slide, text, left, top, width, height, font_size=12, bold=False, color=BODY_COLOR, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    return txBox

def add_formula_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = FORMULA_BG
    shape.line.color.rgb = FORMULA_BORDER
    shape.line.width = Pt(2)
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = FORMULA_TEXT
    p.font.name = 'Consolas'
    return shape

def add_paragraph_block(slide, paragraphs, start_top=TEXT_START_TOP):
    y = start_top
    for text, is_label, is_formula in paragraphs:
        if is_formula:
            add_formula_box(slide, text, LEFT, y, Inches(11.133), FORMULA_HEIGHT)
            y += FORMULA_HEIGHT + SPACING
        elif is_label:
            add_textbox(slide, text, LEFT, y, Inches(11.133), BOLD_HEIGHT,
                        font_size=14, bold=True, color=LABEL_COLOR)
            y += BOLD_HEIGHT + SPACING
        else:
            # Для довгих текстів даємо більше висоти
            est_height = Inches(0.5) if len(text) > 300 else PARAGRAPH_HEIGHT
            add_textbox(slide, text, LEFT, y, Inches(11.133), est_height,
                        font_size=12, color=BODY_COLOR)
            y += est_height + SPACING

# ============================================================
# СЛАЙД 1 — ТИТУЛЬНИЙ
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, BG_COLOR)

# Бейдж "Проєкт із фізики"
badge_s1 = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.0), Inches(3.333), Inches(0.5)
)
badge_s1.fill.solid()
badge_s1.fill.fore_color.rgb = PHENOM_BG
badge_s1.line.fill.background()
p = badge_s1.text_frame.paragraphs[0]
p.text = "Проєкт із фізики"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT
p.alignment = PP_ALIGN.CENTER

# Головна назва
add_textbox(slide1, "Фізичний чек-лист мого дня",
            Inches(1.5), Inches(1.8), Inches(10.333), Inches(0.8),
            font_size=42, bold=True, color=TEXT_COLOR).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Підзаголовок
add_textbox(slide1, "МКТ та термодинаміка",
            Inches(1.5), Inches(2.7), Inches(10.333), Inches(0.5),
            font_size=24, color=SUBTITLE_COLOR).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Роздільник
div = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(3.4), Inches(0.933), Inches(0.05)
)
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0xE8, 0xC5, 0xC5)
div.line.fill.background()

# Автор
add_textbox(slide1, "Підготувала: Вікторія Ляховська",
            Inches(1.5), Inches(3.7), Inches(10.333), Inches(0.4),
            font_size=18, color=SUBTITLE_COLOR).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Девіз
add_textbox(slide1, "Мій день — це суцільний рух молекул, фазові переходи, теплові процеси та робота внутрішньої енергії. Я просто навчилася їх помічати й пояснювати.",
            Inches(2.0), Inches(4.3), Inches(9.333), Inches(0.9),
            font_size=14, color=MOTTO_COLOR).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Додатковий текст
add_textbox(slide1, "Сьогодні я прослідкую за кожним явищем, яке мене супроводжує: від запаху кави зранку до охолодження вечірнього чаю. У цьому проєкті — справжній фізичний щоденник звичайного шкільного дня.",
            Inches(1.5), Inches(5.4), Inches(10.333), Inches(0.7),
            font_size=12, color=EXTRA_COLOR).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# ДАНІ ДЛЯ СЛАЙДІВ 2–10
# ============================================================
all_slides = [
    # Слайд 2
    {
        "badge": "🌅 Ранок",
        "badge_bg": MORNING_BG,
        "badge_color": MORNING_TEXT,
        "title": "Як до мене прийшов запах кави",
        "phenomenon": "Дифузія",
        "paragraphs": [
            ("Я прокинулась і ще лежу в ліжку, але вже чую, як мама порається на кухні. Разом зі звуками до моєї кімнати долинає теплий, глибокий аромат свіжозвареної кави.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Кава містить сотні летких ароматичних речовин (ефірні олії, альдегіди, піразини), які легко випаровуються з гарячої поверхні. У газоподібному стані ці молекули рухаються хаотично з величезними швидкостями — сотні метрів за секунду. Вони стикаються з молекулами азоту та кисню й поступово проникають у простір між ними.", False, False),
            ("Рушійна сила дифузії — різниця концентрацій. На кухні концентрація висока, у спальні — майже нульова. Згідно із законом Фіка, потік речовини пропорційний градієнту концентрації. Швидкість дифузії зростає з температурою, адже Ek = (3/2)kT, тому гаряча кава «пахне» набагато інтенсивніше.", False, False),
            ("Ek = (3/2)kT  |  k = 1,38×10⁻²³ Дж/К  |  T ≈ 373 K → Ek ≈ 7,7×10⁻²¹ Дж", False, True),
        ]
    },
    # Слайд 3
    {
        "badge": "🌅 Ранок",
        "badge_bg": MORNING_BG,
        "badge_color": MORNING_TEXT,
        "title": "Коли я бризкаю улюбленими парфумами",
        "phenomenon": "Випаровування, охолодження під час випаровування, дифузія",
        "paragraphs": [
            ("Перед виходом до школи я роблю один пшик на запʼясток — шкіра стає вологою та приємно прохолодною, а за мить кімната наповнюється квітковим шлейфом.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Туалетна вода — це суміш летких ароматичних речовин у етиловому спирті. Спирт дуже швидко випаровується. Щоб випаруватися, молекули спирту повинні подолати сили міжмолекулярного притягання (водневі звʼязки та ван-дер-ваальсові сили). Це можуть зробити лише найшвидші молекули, які мають кінетичну енергію, більшу за роботу виходу. Покидаючи рідину, вони забирають енергію з собою — середня кінетична енергія молекул, що залишилися, зменшується, температура знижується, і я відчуваю прохолоду.", False, False),
            ("Q = L·m  |  L (спирт) ≈ 840 кДж/кг  |  адіабатичне випаровування", False, True),
            ("💡 Парфуми наносять на точки пульсу (запʼястя, шия, за вухами) — там шкіра тепліша через близькість судин, випаровування інтенсивніше, і аромат розкривається яскравіше.", False, False),
            ("А далі молекули аромату дифундують у повітрі, поширюючись далеко за межі запʼястка — і ось уже в передпокої чути квітковий шлейф.", False, False),
        ]
    },
    # Слайд 4
    {
        "badge": "🌅 Ранок",
        "badge_bg": MORNING_BG,
        "badge_color": MORNING_TEXT,
        "title": "Прасування шкільної блузки",
        "phenomenon": "Теплопровідність, пароутворення, конденсація, фазовий перехід",
        "paragraphs": [
            ("Я вмикаю праску, чекаю, поки нагріється, і веду по зімʼятій тканині. За кілька секунд блузка стає ідеально рівною, а від тканини йде легка пара.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Металева підошва нагрівається завдяки електричному струму. Метали мають високу теплопровідність завдяки вільним електронам: коефіцієнт теплопровідності алюмінію ≈ 200 Вт/(м·К). Внутрішня енергія металу зростає, коливання іонів кристалічної ґратки посилюються.", False, False),
            ("Коли гаряча підошва торкається вологої тканини, вода поглинає енергію. Коли її стає достатньо для подолання міжмолекулярних звʼязків, вода перетворюється на пару. Питома теплота пароутворення води — 2260 кДж/кг. Пара розширюється, її тиск розпрямляє волокна тканини. Виходячи назовні й контактуючи з холоднішим повітрям, пара конденсується — я бачу легкий білий туман. Це замкнений цикл фазових переходів.", False, False),
            ("Закон Фур'є: Q/Δt = −k·A·(ΔT/Δx)  |  k (алюміній) ≈ 200 Вт/(м·К)", False, True),
        ]
    },
    # Слайд 5
    {
        "badge": "🌅 Ранок",
        "badge_bg": MORNING_BG,
        "badge_color": MORNING_TEXT,
        "title": "Чайна ложка обпікає пальці",
        "phenomenon": "Теплопровідність",
        "paragraphs": [
            ("Я наливаю окріп у горнятко, кидаю цукор і починаю розмішувати. Уже за кілька секунд металева ложка стає гарячою — доводиться швидко витягти її.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Гаряча вода (≈100 °C) має велику внутрішню енергію: молекули рухаються дуже швидко. Коли ложка занурена, молекули води стикаються з іонами кристалічної ґратки металу й передають їм кінетичну енергію. У металах є вільні електрони, які дуже швидко поширюють отриману енергію по всьому обʼєму — саме тому метали є чудовими провідниками тепла.", False, False),
            ("k (сталь) ≈ 50 Вт/(м·К)  |  набагато більше, ніж у дерева (~0,15) чи пластику (~0,2)", False, True),
            ("Теплообмін триватиме, поки температура води й ложки не вирівняється — тоді настане теплова рівновага. Мої пальці, торкаючись гарячого металу, отримують тепло через ту саму теплопровідність — терморецептори шкіри реагують і посилають сигнал у мозок.", False, False),
        ]
    },
    # Слайд 6
    {
        "badge": "☀️ День у школі",
        "badge_bg": DAY_BG,
        "badge_color": DAY_TEXT,
        "title": "Після фізкультури: холодок у коридорі",
        "phenomenon": "Випаровування з поглинанням тепла",
        "paragraphs": [
            ("Ми грали у волейбол, футболка намокла від поту. Коли я вийшла в коридор, мене одразу пройняв холодок — стало навіть мерзлякувато.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Піт на 99% — це вода. Вона починає випаровуватися. Щоб покинути рідину, молекула повинна подолати водневі звʼязки, а це потребує енергії. Питома теплота пароутворення води при 37 °C ≈ 2400 кДж/кг. Для випаровування лише 1 грама поту потрібно ≈2400 Дж — цю енергію молекули забирають зі шкіри. Шкіра охолоджується, кров у капілярах охолоджується — я відчуваю прохолоду.", False, False),
            ("💡 На протязі в коридорі випаровування пришвидшується: потік повітря забирає молекули пари від шкіри, запобігаючи насиченню. Це вимушена конвекція. При високій вологості випаровування сповільнюється — тому в спеку з високою вологістю нам важче.", False, False),
            ("Це природний механізм терморегуляції тіла!", False, True),
        ]
    },
    # Слайд 7
    {
        "badge": "☀️ День у школі",
        "badge_bg": DAY_BG,
        "badge_color": DAY_TEXT,
        "title": "Пил у сонячному промені",
        "phenomenon": "Броунівський рух",
        "paragraphs": [
            ("Сонце пробивається крізь вікно, і в золотому промені я бачу, як танцюють дрібні порошинки. Вони не падають униз, а хаотично ворушаться, ніби живі.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Броунівський рух — безперервний хаотичний рух дрібних частинок, завислих у газі чи рідині, спричинений ударами молекул середовища. Молекули азоту та кисню рухаються зі швидкостями ≈500 м/с (при 20 °C) і безперервно стикаються з порошинкою. Кількість ударів із різних боків у кожен момент різна — виникають флуктуації тиску, і частинка отримує випадкові поштовхи.", False, False),
            ("⟨x²⟩ = 2Dt  |  D — коефіцієнт дифузії (залежить від T, вʼязкості, розміру частинки)", False, True),
            ("📜 Відкритий Робертом Броуном у 1827 р. У 1905 р. Альберт Ейнштейн розробив математичну теорію — це стало вагомим доказом існування атомів і молекул. Дивлячись на танцюючий пил, я на власні очі бачу наслідок теплового руху молекул.", False, False),
        ]
    },
    # Слайд 8
    {
        "badge": "🏠 День після школи",
        "badge_bg": DAY_BG,
        "badge_color": DAY_TEXT,
        "title": "Розігріваю обід: суп закипає",
        "phenomenon": "Теплопередача, конвекція, кипіння, пароутворення",
        "paragraphs": [
            ("Я ставлю каструльку на плиту. Вогонь нагріває дно, рідина починає рухатись, а згодом у всій товщі зʼявляються бульбашки — суп закипів.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Спочатку тепло від полумʼя передається дну через теплопровідність. Нагріта рідина розширюється, її густина зменшується, вона підіймається — виникає природна конвекція. Холодніша рідина опускається, і так перемішується весь обʼєм.", False, False),
            ("Коли температура сягає ≈100 °C (за нормального атм. тиску 101325 Па), починається кипіння — пароутворення в усьому обʼємі. На мікронерівностях дна утворюються бульбашки насиченої пари. Тиск пари всередині має зрівноважити атмосферний + гідростатичний тиск. Бульбашки спливають і лопаються. Питома теплота пароутворення води — 2260 кДж/кг. Під час кипіння температура стала — уся енергія йде на розрив міжмолекулярних звʼязків.", False, False),
            ("🏔️ На високогірʼї атмосферний тиск нижчий → температура кипіння знижується → приготування їжі займає більше часу.", False, False),
        ]
    },
    # Слайд 9
    {
        "badge": "🌙 Вечір",
        "badge_bg": EVENING_BG,
        "badge_color": EVENING_TEXT,
        "title": "Ноутбук, домашнє завдання і гарячий корпус",
        "phenomenon": "Зміна внутрішньої енергії, закон Джоуля — Ленца, вимушена конвекція",
        "paragraphs": [
            ("Я сідаю за уроки, вмикаю ноутбук. Через пів години гуде кулер, а корпус знизу став гарячим — навіть крізь джинси відчутно.", False, False),
            ("🔬 Пояснення:", True, False),
            ("Усередині ноутбука — мільйони транзисторів. Під час їхнього перемикання через мікросхеми тече струм. Закон Джоуля — Ленца: Q = I²Rt. Електрична енергія перетворюється на внутрішню, температура процесора зростає (може сягати 100 °C!).", False, False),
            ("Q = I²Rt  |  Потужність процесора: 15–45 Вт  |  Закон Ньютона — Ріхмана: Q = h·A·(T_proc – T_air)", False, True),
            ("Система охолодження: теплові трубки (рідина випаровується біля процесора й конденсується біля радіатора) + вентилятор (вимушена конвекція). Частина тепла все одно передається корпусу через теплопровідність. Якщо система не впорається — спрацює тротлінг або ноутбук вимкнеться.", False, False),
        ]
    },
    # Слайд 10
    {
        "badge": "🌙 Вечір",
        "badge_bg": EVENING_BG,
        "badge_color": EVENING_TEXT,
        "title": "Горнятко чаю та три шляхи тепловтрат",
        "phenomenon": "Теплопровідність, конвекція, теплове випромінювання, теплова рівновага",
        "paragraphs": [
            ("Я вмостилася з горнятком гарячого чаю. Він поступово холоне — працюють одразу три механізми теплопередачі.", False, False),
            ("🔬 Три механізми в одному горнятку:", True, False),
            ("1. Теплопровідність — тепло від чаю передається через стінки керамічного горнятка (k ≈ 1–2 Вт/(м·К)). Моя долоня гріється саме через теплопровідність.", False, False),
            ("2. Конвекція — повітря біля гарячого горнятка нагрівається, розширюється, стає менш густим і підіймається вгору, забираючи тепло.", False, False),
            ("3. Теплове випромінювання — горнятко випромінює інфрачервоні хвилі. Закон Стефана — Больцмана: P = εσAT⁴. Для кераміки ε ≈ 0,9. При T ≈ 350 K і A ≈ 0,03 м² потужність випромінювання ≈ 20 Вт — суттєвий канал тепловтрат!", False, False),
            ("P = εσAT⁴  |  σ = 5,67×10⁻⁸ Вт/(м²·К⁴)  |  ε (кераміка) ≈ 0,9", False, True),
            ("Чай холоне, внутрішня енергія зменшується. Коли температури зрівняються — настане теплова рівновага (другий закон термодинаміки).", False, False),
        ]
    },
]

# ===== Будуємо слайди 2–10 =====
for sdata in all_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_card(slide)
    add_badge(slide, sdata["badge"], sdata["badge_bg"], sdata["badge_color"])
    add_title(slide, sdata["title"])
    add_phenomenon(slide, sdata["phenomenon"])
    add_paragraph_block(slide, sdata["paragraphs"])

# ============================================================
# СЛАЙД 11 — ВИСНОВКИ
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, BG_COLOR)
add_card(slide11)

add_badge(slide11, "📋 Підсумки", PHENOM_BG, ACCENT, LEFT, BADGE_TOP)
add_title(slide11, "Висновки: мій фізичний день")

conclusions = [
    "Дифузія — запах кави та парфумів, закон Фіка, градієнт концентрації.",
    "Випаровування — парфуми та піт; Q = Lm, адіабатичне охолодження шкіри.",
    "Теплопровідність — ложка в чаї, праска; закон Фур'є, вільні електрони металів.",
    "Фазові переходи — прасування: рідина → пара → рідина, цикл із поглинанням і виділенням теплоти.",
    "Броунівський рух — пил у промені; ⟨x²⟩ = 2Dt, доказ існування молекул.",
    "Конвекція — суп у каструлі, охолодження ноутбука; природна й вимушена.",
    "Кипіння — бульбашки пари; тиск насиченої пари = зовнішній тиск, T = const.",
    "Закон Джоуля — Ленца — ноутбук; Q = I²Rt, електроенергія → внутрішня енергія.",
    "Теплове випромінювання — горнятко; P = εσAT⁴, інфрачервоні хвилі.",
    "Теплова рівновага — другий закон термодинаміки; тепло мимовільно переходить від гарячого до холодного.",
]

y = Inches(2.5)
for c in conclusions:
    add_textbox(slide11, "• " + c, LEFT, y, Inches(11.133), Inches(0.28), font_size=12, color=BODY_COLOR)
    y += Inches(0.28)

y += Inches(0.15)
# Фінальний блок
final_shape = slide11.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, LEFT, y, Inches(11.133), Inches(1.2)
)
final_shape.fill.solid()
final_shape.fill.fore_color.rgb = FORMULA_BG
final_shape.line.color.rgb = FORMULA_BORDER
final_shape.line.width = Pt(2)
tf = final_shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("💭 Мій особистий підсумок: Раніше я думала, що фізика — це формули на дошці та задачі в підручнику. "
          "Тепер я знаю: молекулярна фізика і термодинаміка — це справжня мова мого звичайного дня. "
          "Запах кави, прохолода від парфумів, тепло ноутбука, праска, чайна ложка, навіть пил у сонячному світлі — "
          "усе це розповідає історію про рух молекул, енергію та фазові переходи. "
          "Я більше ніколи не подивлюся на звичайний ранок так, як раніше. "
          "І це робить кожен мій день трішечки чарівним — чарівним і науковим водночас.")
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = LABEL_COLOR

# ===== Збереження =====
output_file = "Фізичний_чек_лист_Вікторія_Ляховська.pptx"
prs.save(output_file)
print(f"✅ Презентацію збережено: {output_file}")